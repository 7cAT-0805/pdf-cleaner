from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
import os

# 在導入 torch 之前設置環境變數，以最大限度減少記憶體使用
os.environ['OMP_NUM_THREADS'] = '1'
os.environ['MKL_NUM_THREADS'] = '1'

import easyocr
import cv2
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import fitz
import tempfile
import uuid
import traceback
import torch
import time
import threading
import base64
import io
import gc

# 限制 PyTorch 緒數以減少記憶體佔用
torch.set_num_threads(1)

app = Flask(__name__)
CORS(app)

# 全域變數：儲存即時處理狀態
processing_status = {
    "current_page": 0,
    "total_pages": 0,
    "preview_image": None,
    "is_processing": False
}

# 移除全域 reader 初始化，改為延遲載入函數
_reader = None
_reader_lock = threading.Lock()

def get_reader():
    global _reader
    with _reader_lock:
        if _reader is None:
            print("[*] 正在延遲初始化 EasyOCR 引擎...")
            model_path = os.path.join(os.path.dirname(__file__), 'models')
            # 增加 recognizer=True (預設), 但確保我們限制了模型搜尋路徑
            has_gpu = torch.cuda.is_available()
            _reader = easyocr.Reader(['ch_tra', 'en'], gpu=has_gpu, model_storage_directory=model_path, download_enabled=False)
            print("[*] EasyOCR 引擎初始化完成")
        return _reader

def release_reader():
    global _reader
    with _reader_lock:
        if _reader is not None:
            print("[*] 釋放 EasyOCR 引擎以節省記憶體...")
            del _reader
            _reader = None
            gc.collect()
            if torch.cuda.is_available():
                torch.cuda.empty_cache()

def get_actual_text_color(original_img, bbox):
    try:
        coords = np.array(bbox).astype(np.int32)
        xmin, ymin = np.min(coords, axis=0)
        xmax, ymax = np.max(coords, axis=0)
        
        crop = original_img[int(max(0,ymin+1)):int(min(original_img.shape[0],ymax-1)), 
                            int(max(0,xmin+1)):int(min(original_img.shape[1],xmax-1))]
        
        if crop.size == 0: return RGBColor(0, 0, 0), False
        
        pixels = crop.reshape(-1, 3).astype(np.float32)
        brightness = 0.299 * pixels[:, 2] + 0.587 * pixels[:, 1] + 0.114 * pixels[:, 0]
        avg_b = np.mean(brightness)
        
        if avg_b > 127:
            target_mask = brightness < np.percentile(brightness, 25)
        else:
            target_mask = brightness > np.percentile(brightness, 75)
            
        selected_pixels = pixels[target_mask]
        if selected_pixels.size == 0:
            return RGBColor(0, 0, 0), False
            
        color = np.mean(selected_pixels, axis=0)
        is_bold = (selected_pixels.shape[0] / pixels.shape[0]) > 0.35
        
        return RGBColor(int(color[2]), int(color[1]), int(color[0])), is_bold
    except:
        return RGBColor(0, 0, 0), False

def inpaint_text(img, bboxes):
    mask = np.zeros(img.shape[:2], dtype=np.uint8)
    for bbox in bboxes:
        pts = np.array(bbox).astype(np.int32)
        cv2.fillPoly(mask, [pts], 255)
    kernel = np.ones((9,9), np.uint8)
    mask = cv2.dilate(mask, kernel, iterations=1)
    return cv2.inpaint(img, mask, 3, cv2.INPAINT_TELEA)

@app.route('/status', methods=['GET'])
def get_status():
    """返回當前處理狀態和預覽圖"""
    global processing_status
    return jsonify({
        "current_page": processing_status["current_page"],
        "total_pages": processing_status["total_pages"],
        "preview": processing_status["preview_image"],
        "is_processing": processing_status["is_processing"]
    })

def calculate_rotation_angle(bbox):
    """
    從 EasyOCR 的 bbox 計算文字旋轉角度（度）
    bbox: [[x1,y1], [x2,y2], [x3,y3], [x4,y4]]
    返回：-180 到 180 度之間的角度
    """
    try:
        # 使用第一個點到第二個點的向量來計算角度
        x1, y1 = bbox[0]
        x2, y2 = bbox[1]
        
        # 計算角度（弧度轉角度）
        angle_rad = np.arctan2(y2 - y1, x2 - x1)
        angle_deg = np.degrees(angle_rad)
        
        # 將角度限制在 -45 到 45 度之間（PPT 相容性）
        # 如果角度太極端，可能是垂直文字，我們暫時不處理
        if abs(angle_deg) > 45:
            return 0
            
        return angle_deg  # 直接返回計算出的角度
    except:
        return 0

@app.route('/')
def index():
    return "PDF to PPT Converter API is running", 200

@app.route('/health')
def health_check():
    return "OK", 200

@app.route('/convert', methods=['POST'])
def convert_pdf_to_ppt():
    global processing_status
    try:
        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400
        pdf_file = request.files['file']
        # Use a temporary directory for intermediate files
        temp_dir = tempfile.mkdtemp()
        pdf_path = os.path.join(temp_dir, f"{uuid.uuid4()}.pdf")
        pdf_file.save(pdf_path)
        doc = fitz.open(pdf_path)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        processing_status["total_pages"] = doc.page_count
        processing_status["is_processing"] = True
        
        # 獲取 Reader，如果沒有則初始化
        reader = get_reader()
        
        # Process each page
        for pno in range(doc.page_count):
            processing_status["current_page"] = pno + 1
            print(f"[*] AI 深度分析: {pno+1}/{doc.page_count}")
            page = doc.load_page(pno)
            pix = page.get_pixmap(dpi=150) # 降低 DPI 從 300 到 150 以節省記憶體
            
            # 直接從 pixmap 取得 numpy 陣列，避免重複的 PNG 編解碼
            raw_img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.h, pix.w, pix.n)
            if pix.n == 4:
                raw_img = cv2.cvtColor(raw_img, cv2.COLOR_RGBA2RGB)
                raw_img = cv2.cvtColor(raw_img, cv2.COLOR_RGB2BGR)
            elif pix.n == 1:
                raw_img = cv2.cvtColor(raw_img, cv2.COLOR_GRAY2BGR)
            else:
                raw_img = cv2.cvtColor(raw_img, cv2.COLOR_RGB2BGR)
                
            # OCR (直接傳入 numpy 陣列)
            results = reader.readtext(raw_img, contrast_ths=0.1, width_ths=1.0, y_ths=0.1)
            
            # Inpaint
            bboxes = [res[0] for res in results]
            clean_img = inpaint_text(raw_img, bboxes)
            # Preview (small version) for frontend
            preview = cv2.resize(clean_img, (600, int(600 * clean_img.shape[0] / clean_img.shape[1])))
            _, buffer = cv2.imencode('.jpg', preview, [cv2.IMWRITE_JPEG_QUALITY, 70])
            processing_status["preview_image"] = base64.b64encode(buffer).decode('utf-8')
            # Add slide image directly from numpy array
            _, img_buf = cv2.imencode('.png', clean_img)
            img_bytes = img_buf.tobytes()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(io.BytesIO(img_bytes), 0, 0, prs.slide_width, prs.slide_height)
            
            # 準備文字框
            img_w, img_h = pix.width, pix.height
            
            # Add text boxes
            for (bbox, text, prob) in results:
                if prob < 0.15:
                    continue
                coords = np.array(bbox)
                xmin, ymin = np.min(coords, axis=0)
                xmax, ymax = np.max(coords, axis=0)
                final_x = int((xmin / img_w) * prs.slide_width)
                final_y = int((ymin / img_h) * prs.slide_height)
                final_w = int(((xmax - xmin) / img_w) * prs.slide_width * 1.1)
                final_h = int(((ymax - ymin) / img_h) * prs.slide_height)
                rotation_angle = calculate_rotation_angle(bbox)
                text_color, is_bold = get_actual_text_color(raw_img, bbox)
                txBox = slide.shapes.add_textbox(final_x, final_y, final_w, final_h)
                txBox.margin_left = txBox.margin_right = txBox.margin_top = txBox.margin_bottom = 0
                txBox.rotation = rotation_angle
                tf = txBox.text_frame
                tf.text = text
                tf.word_wrap = False
                p = tf.paragraphs[0]
                p.font.color.rgb = text_color
                p.font.bold = is_bold
                p.font.name = "Microsoft JhengHei"
                fontsize = int((final_h / 12700) * 0.72)
                p.font.size = Pt(min(max(10, fontsize), 72))
            
            # 清理該頁資源
            del raw_img, clean_img, pix, page, results, bboxes
            gc.collect()

        doc.close()
        output_pptx = os.path.join(temp_dir, "output.pptx")
        prs.save(output_pptx)
        # Reset status and clear preview
        processing_status["is_processing"] = False
        processing_status["preview_image"] = None
        
        # 處理完成後釋放 EasyOCR 以節省記憶體
        release_reader()

        # Send file and clean up temp directory after response
        def cleanup_and_send():
            try:
                return send_file(output_pptx, as_attachment=True, download_name="SMART_PRO_v2.3.1.pptx")
            finally:
                # Remove temporary files
                try:
                    os.remove(output_pptx)
                    os.remove(pdf_path)
                    os.rmdir(temp_dir)
                except Exception as e:
                    print(f"Cleanup error: {e}")
        return cleanup_and_send()
    except Exception as e:
        processing_status["is_processing"] = False
        processing_status["preview_image"] = None
        release_reader()
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    app.run(port=5000)
